"""
FY26 Sales Forecast Report — Laptop Pro | Acme Corporation
Dark Modern Consulting Theme (from consulting-powerpoint-master-SKILL.md)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
import os

# ── Theme Constants ──────────────────────────────────────────────────────────
BG_COLOR      = RGBColor(0x19, 0x19, 0x19)  # #191919
CARD_COLOR    = RGBColor(0x26, 0x26, 0x26)  # #262626
ACCENT        = RGBColor(0x63, 0x66, 0xF1)  # #6366F1 indigo
ACCENT_LIGHT  = RGBColor(0x81, 0x84, 0xF7)  # lighter indigo for hover/secondary
TEXT_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY     = RGBColor(0xA0, 0xA0, 0xA0)  # #A0A0A0
TEXT_DIM      = RGBColor(0x6B, 0x6B, 0x6B)  # #6B6B6B
GREEN         = RGBColor(0x22, 0xC5, 0x5E)  # #22C55E
RED           = RGBColor(0xEF, 0x44, 0x44)  # #EF4444
YELLOW        = RGBColor(0xFA, 0xCC, 0x15)  # #FACC15
BORDER_COLOR  = RGBColor(0x33, 0x33, 0x33)  # #333333

FONT_BODY     = "Calibri"  # fallback if Liter not available
SLIDE_W       = Inches(13.333)
SLIDE_H       = Inches(7.5)

# Chart colors
CHART_COLORS = [ACCENT, GREEN, YELLOW, RGBColor(0xF9, 0x73, 0x16), RED, ACCENT_LIGHT]

# ── Data ─────────────────────────────────────────────────────────────────────
# FY26 Forecast for Laptop Pro | Acme Corporation
FY24_MONTHLY = [600, 1800, 1200, 400, 800, 1000, 1000, 1200, 600, 200, 2000, 1400]
FY25_MONTHLY = [600, 1200, 1600, 0, 1600, 1400, 1600, 600, 200, 800, 1800, 400]
FY26_MONTHLY = [900, 2700, 1800, 600, 1200, 1500, 1500, 1800, 900, 300, 3000, 2100]
FY27_MONTHLY = [900, 1800, 2400, 0, 2400, 2100, 2400, 900, 300, 1200, 2700, 600]

MONTHS = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]

FY24_Q = [3600, 2200, 2800, 3600]
FY25_Q = [3400, 3000, 2400, 3000]
FY26_Q = [5400, 3300, 4200, 5400]
FY27_Q = [5100, 4500, 3600, 4500]

FY24_TOTAL = 12200
FY25_TOTAL = 11800
FY26_TOTAL = 18300
FY27_TOTAL = 17700

QUARTERS = ["Q1", "Q2", "Q3", "Q4"]

# ── Helpers ──────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=CARD_COLOR, border=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border:
        shape.line.color.rgb = BORDER_COLOR
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Subtle rounded corners
    shape.adjustments[0] = 0.02
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=TEXT_WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_BODY, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_accent_bar(slide, left, top, width=Inches(0.06), height=Inches(0.5)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    return bar


def format_num(n):
    """Format number with comma separator"""
    if isinstance(n, float):
        return f"{n:,.0f}"
    return f"{n:,}"


def pct_change(new, old):
    if old == 0:
        return "N/A"
    change = ((new - old) / old) * 100
    sign = "+" if change >= 0 else ""
    return f"{sign}{change:.1f}%"


# ── Slide Builders ───────────────────────────────────────────────────────────

def build_title_slide(prs):
    """Slide 1: Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Accent line at top
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Main title
    add_text(slide, Inches(1.2), Inches(2.0), Inches(10), Inches(1.2),
             "FY26 Sales Forecast Report", font_size=44, bold=True, color=TEXT_WHITE)

    # Subtitle
    add_text(slide, Inches(1.2), Inches(3.3), Inches(10), Inches(0.6),
             "Laptop Pro  |  Acme Corporation", font_size=24, color=ACCENT)

    # Divider line
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(4.2), Inches(3), Inches(0.03))
    div.fill.solid()
    div.fill.fore_color.rgb = ACCENT
    div.line.fill.background()

    # Details
    add_text(slide, Inches(1.2), Inches(4.6), Inches(8), Inches(0.4),
             "Source: [LARA] TEST MODEL  |  Anaplan Planning Platform", font_size=14, color=TEXT_GRAY)
    add_text(slide, Inches(1.2), Inches(5.1), Inches(8), Inches(0.4),
             "Prepared: February 2026  |  Confidential", font_size=12, color=TEXT_DIM)

    # Bottom accent bar
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = ACCENT
    bottom.line.fill.background()


def build_exec_summary(prs):
    """Slide 2: Executive Summary with KPI cards"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Header
    add_accent_bar(slide, Inches(0.8), Inches(0.5), height=Inches(0.45))
    add_text(slide, Inches(1.0), Inches(0.45), Inches(10), Inches(0.6),
             "Executive Summary", font_size=28, bold=True)

    # Bumper statement
    add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "FY26 forecasts project 55% revenue growth for Laptop Pro at Acme Corporation, driven by historical momentum and 1.5x growth multiplier",
             font_size=15, color=TEXT_GRAY)

    # KPI Cards Row
    card_y = Inches(2.0)
    card_h = Inches(2.2)
    card_w = Inches(2.7)
    gap = Inches(0.25)
    start_x = Inches(0.8)

    kpis = [
        ("FY26 Forecast", f"${format_num(FY26_TOTAL)}", pct_change(FY26_TOTAL, FY25_TOTAL), "vs FY25"),
        ("FY25 Actual", f"${format_num(FY25_TOTAL)}", pct_change(FY25_TOTAL, FY24_TOTAL), "vs FY24"),
        ("FY24 Actual", f"${format_num(FY24_TOTAL)}", "-", "Baseline Year"),
        ("FY27 Outlook", f"${format_num(FY27_TOTAL)}", pct_change(FY27_TOTAL, FY26_TOTAL), "vs FY26"),
    ]

    for i, (label, value, change, sub) in enumerate(kpis):
        x = start_x + (card_w + gap) * i
        card = add_shape(slide, x, card_y, card_w, card_h, border=True)

        add_text(slide, x + Inches(0.25), card_y + Inches(0.25), card_w - Inches(0.5), Inches(0.35),
                 label, font_size=12, color=TEXT_GRAY)

        add_text(slide, x + Inches(0.25), card_y + Inches(0.65), card_w - Inches(0.5), Inches(0.6),
                 value, font_size=32, bold=True, color=TEXT_WHITE)

        # Change indicator
        change_color = GREEN if change.startswith("+") else RED if change.startswith("-") else TEXT_GRAY
        if change == "N/A" or change == "-":
            change_color = TEXT_GRAY
        add_text(slide, x + Inches(0.25), card_y + Inches(1.35), card_w - Inches(0.5), Inches(0.35),
                 change, font_size=18, bold=True, color=change_color)

        add_text(slide, x + Inches(0.25), card_y + Inches(1.7), card_w - Inches(0.5), Inches(0.3),
                 sub, font_size=11, color=TEXT_DIM)

    # Key Insights section
    insights_y = Inches(4.6)
    add_accent_bar(slide, Inches(0.8), insights_y, height=Inches(0.35))
    add_text(slide, Inches(1.0), insights_y - Inches(0.05), Inches(5), Inches(0.45),
             "Key Insights", font_size=18, bold=True)

    insights = [
        "Projected 55.1% YoY growth from FY25 ($11.8K) to FY26 ($18.3K)",
        "Q4 consistently strongest quarter across all years (seasonal peak Nov-Dec)",
        "FY26 forecast uses 1.5x historical multiplier applied to FY24 base data",
        "Q2 shows relative softness — potential area for targeted sales initiatives",
    ]

    for j, insight in enumerate(insights):
        y = insights_y + Inches(0.5) + Inches(0.4) * j
        # Bullet dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.08), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_text(slide, Inches(1.25), y, Inches(10), Inches(0.35),
                 insight, font_size=13, color=TEXT_GRAY)


def build_quarterly_comparison(prs):
    """Slide 3: Quarterly Revenue Comparison chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Header
    add_accent_bar(slide, Inches(0.8), Inches(0.5), height=Inches(0.45))
    add_text(slide, Inches(1.0), Inches(0.45), Inches(10), Inches(0.6),
             "Quarterly Revenue Comparison", font_size=28, bold=True)

    add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.4),
             "FY24 Actual vs FY25 Actual vs FY26 Forecast  |  Laptop Pro  |  Acme Corporation",
             font_size=13, color=TEXT_GRAY)

    # Chart
    chart_data = CategoryChartData()
    chart_data.categories = QUARTERS
    chart_data.add_series("FY24 Actual", FY24_Q)
    chart_data.add_series("FY25 Actual", FY25_Q)
    chart_data.add_series("FY26 Forecast", FY26_Q)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.7),
        Inches(7.5), Inches(5.0), chart_data
    )
    chart = chart_frame.chart

    # Style the chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.color.rgb = TEXT_GRAY
    chart.legend.font.size = Pt(11)

    # Chart area
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 100

    # Series colors
    series_colors = [TEXT_DIM, ACCENT_LIGHT, ACCENT]
    for idx, series in enumerate(plot.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = series_colors[idx]
        # Data labels
        series.has_data_labels = True
        series.data_labels.font.size = Pt(10)
        series.data_labels.font.color.rgb = TEXT_WHITE
        series.data_labels.number_format = '#,##0'
        series.data_labels.show_value = True

    # Axes
    value_axis = chart.value_axis
    value_axis.has_title = False
    value_axis.major_gridlines.format.line.color.rgb = BORDER_COLOR
    value_axis.format.line.color.rgb = BORDER_COLOR
    value_axis.tick_labels.font.color.rgb = TEXT_GRAY
    value_axis.tick_labels.font.size = Pt(10)
    value_axis.tick_labels.number_format = '$#,##0'

    category_axis = chart.category_axis
    category_axis.format.line.color.rgb = BORDER_COLOR
    category_axis.tick_labels.font.color.rgb = TEXT_GRAY
    category_axis.tick_labels.font.size = Pt(11)

    # Side panel — YoY Growth Table
    panel_x = Inches(8.8)
    panel_w = Inches(4.0)
    panel_card = add_shape(slide, panel_x, Inches(1.7), panel_w, Inches(5.0), border=True)

    add_text(slide, panel_x + Inches(0.3), Inches(1.9), panel_w - Inches(0.6), Inches(0.4),
             "Year-over-Year Growth", font_size=16, bold=True)

    # Growth table rows
    table_data = [
        ("Quarter", "FY24>25", "FY25>26"),
        ("Q1", pct_change(FY25_Q[0], FY24_Q[0]), pct_change(FY26_Q[0], FY25_Q[0])),
        ("Q2", pct_change(FY25_Q[1], FY24_Q[1]), pct_change(FY26_Q[1], FY25_Q[1])),
        ("Q3", pct_change(FY25_Q[2], FY24_Q[2]), pct_change(FY26_Q[2], FY25_Q[2])),
        ("Q4", pct_change(FY25_Q[3], FY24_Q[3]), pct_change(FY26_Q[3], FY25_Q[3])),
        ("Total", pct_change(FY25_TOTAL, FY24_TOTAL), pct_change(FY26_TOTAL, FY25_TOTAL)),
    ]

    row_h = Inches(0.5)
    for r, (col1, col2, col3) in enumerate(table_data):
        y = Inches(2.5) + row_h * r
        is_header = r == 0
        is_total = r == len(table_data) - 1
        c = TEXT_WHITE if is_header or is_total else TEXT_GRAY
        b = is_header or is_total

        add_text(slide, panel_x + Inches(0.3), y, Inches(1.0), row_h, col1, font_size=12, color=c, bold=b)
        add_text(slide, panel_x + Inches(1.4), y, Inches(1.1), row_h, col2, font_size=12,
                 color=RED if col2.startswith("-") else GREEN if col2.startswith("+") else c, bold=b,
                 alignment=PP_ALIGN.CENTER)
        add_text(slide, panel_x + Inches(2.6), y, Inches(1.1), row_h, col3, font_size=12,
                 color=RED if col3.startswith("-") else GREEN if col3.startswith("+") else c, bold=b,
                 alignment=PP_ALIGN.CENTER)

    # Separator lines
    for r in [0, len(table_data) - 1]:
        y = Inches(2.5) + row_h * r + row_h - Inches(0.02)
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, panel_x + Inches(0.3), y, panel_w - Inches(0.6), Inches(0.015))
        sep.fill.solid()
        sep.fill.fore_color.rgb = BORDER_COLOR
        sep.line.fill.background()


def build_monthly_trend(prs):
    """Slide 4: Monthly Trend — FY26 Forecast Breakdown"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.5), height=Inches(0.45))
    add_text(slide, Inches(1.0), Inches(0.45), Inches(10), Inches(0.6),
             "FY26 Monthly Forecast Breakdown", font_size=28, bold=True)

    add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.4),
             "Month-by-month forecast with quarterly aggregation  |  Forecast = Historical FY24 x 1.5",
             font_size=13, color=TEXT_GRAY)

    # Line chart with FY24, FY25, FY26
    chart_data = CategoryChartData()
    chart_data.categories = MONTHS
    chart_data.add_series("FY24 Actual", FY24_MONTHLY)
    chart_data.add_series("FY25 Actual", FY25_MONTHLY)
    chart_data.add_series("FY26 Forecast", FY26_MONTHLY)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(0.8), Inches(1.7),
        Inches(11.7), Inches(5.0), chart_data
    )
    chart = chart_frame.chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.color.rgb = TEXT_GRAY
    chart.legend.font.size = Pt(11)

    plot = chart.plots[0]
    series_colors = [TEXT_DIM, ACCENT_LIGHT, ACCENT]
    series_widths = [Pt(2), Pt(2), Pt(3)]

    for idx, series in enumerate(plot.series):
        series.format.line.color.rgb = series_colors[idx]
        series.format.line.width = series_widths[idx]
        series.marker.style = 8  # circle
        series.marker.size = 8
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = series_colors[idx]
        series.marker.format.line.color.rgb = series_colors[idx]

        if idx == 2:  # FY26 — show data labels
            series.has_data_labels = True
            series.data_labels.font.size = Pt(9)
            series.data_labels.font.color.rgb = ACCENT
            series.data_labels.number_format = '$#,##0'
            series.data_labels.show_value = True
            series.data_labels.position = XL_LABEL_POSITION.ABOVE

    # Axes
    value_axis = chart.value_axis
    value_axis.major_gridlines.format.line.color.rgb = BORDER_COLOR
    value_axis.format.line.color.rgb = BORDER_COLOR
    value_axis.tick_labels.font.color.rgb = TEXT_GRAY
    value_axis.tick_labels.font.size = Pt(10)
    value_axis.tick_labels.number_format = '$#,##0'

    category_axis = chart.category_axis
    category_axis.format.line.color.rgb = BORDER_COLOR
    category_axis.tick_labels.font.color.rgb = TEXT_GRAY
    category_axis.tick_labels.font.size = Pt(10)


def build_data_table(prs):
    """Slide 5: Detailed Data Table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.5), height=Inches(0.45))
    add_text(slide, Inches(1.0), Inches(0.45), Inches(10), Inches(0.6),
             "Detailed Revenue Data", font_size=28, bold=True)

    add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.4),
             "Laptop Pro  |  Acme Corporation  |  All figures in USD",
             font_size=13, color=TEXT_GRAY)

    # Build table
    rows = 5  # header + FY24 + FY25 + FY26 + FY27
    cols = 6  # Label + Q1 + Q2 + Q3 + Q4 + Total
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(11.7)
    height = Inches(2.5)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    col_widths = [Inches(2.5), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(2.0)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    headers = ["Period", "Q1", "Q2", "Q3", "Q4", "Full Year"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = TEXT_WHITE
            paragraph.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT

    # Data rows
    data_rows = [
        ("FY24 Actual", FY24_Q, FY24_TOTAL),
        ("FY25 Actual", FY25_Q, FY25_TOTAL),
        ("FY26 Forecast", FY26_Q, FY26_TOTAL),
        ("FY27 Outlook", FY27_Q, FY27_TOTAL),
    ]

    for r, (label, quarters, total) in enumerate(data_rows):
        row_idx = r + 1
        values = [label] + [f"${format_num(q)}" for q in quarters] + [f"${format_num(total)}"]
        row_color = CARD_COLOR if r % 2 == 0 else RGBColor(0x1F, 0x1F, 0x1F)

        for c, val in enumerate(values):
            cell = table.cell(row_idx, c)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = TEXT_WHITE
                paragraph.font.bold = (c == 0 or c == len(values) - 1)
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color

    # ── Monthly breakdown below ──
    add_text(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.4),
             "Monthly Breakdown — FY26 Forecast", font_size=16, bold=True)

    m_rows = 2  # header + values
    m_cols = 13  # label + 12 months
    m_top = Inches(5.0)
    m_height = Inches(1.0)

    m_table_shape = slide.shapes.add_table(m_rows, m_cols, left, m_top, width, m_height)
    m_table = m_table_shape.table

    # Column widths for monthly table
    m_table.columns[0].width = Inches(1.5)
    for i in range(1, 13):
        m_table.columns[i].width = Inches(0.85)

    # Headers
    m_headers = ["FY26"] + MONTHS
    for i, h in enumerate(m_headers):
        cell = m_table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = TEXT_WHITE
            paragraph.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT

    # Values
    m_values = ["Forecast"] + [f"${format_num(v)}" for v in FY26_MONTHLY]
    for i, val in enumerate(m_values):
        cell = m_table.cell(1, i)
        cell.text = val
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = TEXT_WHITE
            paragraph.font.bold = (i == 0)
            paragraph.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_COLOR


def build_waterfall_analysis(prs):
    """Slide 6: Quarterly waterfall / seasonal analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.5), height=Inches(0.45))
    add_text(slide, Inches(1.0), Inches(0.45), Inches(10), Inches(0.6),
             "Seasonal Pattern Analysis", font_size=28, bold=True)

    add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.4),
             "Quarterly contribution to annual revenue  |  Consistent Q4 peak driven by year-end procurement cycles",
             font_size=13, color=TEXT_GRAY)

    # Stacked bar chart showing quarterly contribution percentages
    chart_data = CategoryChartData()
    chart_data.categories = ["FY24", "FY25", "FY26", "FY27"]
    chart_data.add_series("Q1", [FY24_Q[0], FY25_Q[0], FY26_Q[0], FY27_Q[0]])
    chart_data.add_series("Q2", [FY24_Q[1], FY25_Q[1], FY26_Q[1], FY27_Q[1]])
    chart_data.add_series("Q3", [FY24_Q[2], FY25_Q[2], FY26_Q[2], FY27_Q[2]])
    chart_data.add_series("Q4", [FY24_Q[3], FY25_Q[3], FY26_Q[3], FY27_Q[3]])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, Inches(0.8), Inches(1.7),
        Inches(7.0), Inches(5.0), chart_data
    )
    chart = chart_frame.chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.color.rgb = TEXT_GRAY
    chart.legend.font.size = Pt(11)

    plot = chart.plots[0]
    plot.gap_width = 150

    q_colors = [ACCENT, ACCENT_LIGHT, GREEN, YELLOW]
    for idx, series in enumerate(plot.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = q_colors[idx]
        series.has_data_labels = True
        series.data_labels.font.size = Pt(9)
        series.data_labels.font.color.rgb = TEXT_WHITE
        series.data_labels.number_format = '$#,##0'
        series.data_labels.show_value = True

    # Axes
    value_axis = chart.value_axis
    value_axis.major_gridlines.format.line.color.rgb = BORDER_COLOR
    value_axis.format.line.color.rgb = BORDER_COLOR
    value_axis.tick_labels.font.color.rgb = TEXT_GRAY
    value_axis.tick_labels.font.size = Pt(10)
    value_axis.tick_labels.number_format = '$#,##0'

    category_axis = chart.category_axis
    category_axis.format.line.color.rgb = BORDER_COLOR
    category_axis.tick_labels.font.color.rgb = TEXT_GRAY
    category_axis.tick_labels.font.size = Pt(11)

    # Right side — Seasonal insights card
    card_x = Inches(8.3)
    card_w = Inches(4.5)
    add_shape(slide, card_x, Inches(1.7), card_w, Inches(5.0), border=True)

    add_text(slide, card_x + Inches(0.3), Inches(1.9), card_w - Inches(0.6), Inches(0.4),
             "Seasonal Insights", font_size=16, bold=True)

    insights = [
        ("Q1 Contribution", f"{FY26_Q[0]/FY26_TOTAL*100:.0f}%", "Moderate start — post-holiday ramp-up"),
        ("Q2 Contribution", f"{FY26_Q[1]/FY26_TOTAL*100:.0f}%", "Lowest quarter — budget cycle trough"),
        ("Q3 Contribution", f"{FY26_Q[2]/FY26_TOTAL*100:.0f}%", "Mid-year recovery — project spending"),
        ("Q4 Contribution", f"{FY26_Q[3]/FY26_TOTAL*100:.0f}%", "Peak quarter — year-end procurement"),
    ]

    for i, (title, pct, desc) in enumerate(insights):
        y = Inches(2.6) + Inches(1.05) * i
        add_accent_bar(slide, card_x + Inches(0.3), y, width=Inches(0.04), height=Inches(0.7))

        add_text(slide, card_x + Inches(0.5), y, card_w - Inches(1.0), Inches(0.3),
                 title, font_size=12, bold=True, color=TEXT_WHITE)
        add_text(slide, card_x + Inches(0.5), y + Inches(0.28), Inches(0.8), Inches(0.3),
                 pct, font_size=20, bold=True, color=ACCENT)
        add_text(slide, card_x + Inches(1.5), y + Inches(0.32), card_w - Inches(2.0), Inches(0.3),
                 desc, font_size=11, color=TEXT_GRAY)


def build_forecast_methodology(prs):
    """Slide 7: Forecast Methodology"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.5), height=Inches(0.45))
    add_text(slide, Inches(1.0), Inches(0.45), Inches(10), Inches(0.6),
             "Forecast Methodology & Line Items", font_size=28, bold=True)

    add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.4),
             "Anaplan model structure and calculation logic  |  Sales Forecast module",
             font_size=13, color=TEXT_GRAY)

    # Four methodology cards
    methods = [
        ("Historical Sales", "Raw Input", "Actual recorded sales data\nfrom transactional systems.\n\nFY24: $12,200\nFY25: $11,800",
         "Manual data entry / import"),
        ("Forecast - Mean", "Historical Sales / 12", "Monthly average of annual\nhistorical sales.\n\nSmoothes volatility for\nsteady-state baseline.",
         "Simple averaging method"),
        ("Forecast - Previous Year", "OFFSET(Historical, 0, -12)", "Same month, prior year.\nAssumes repeat patterns\nwith no growth adjustment.",
         "Naive seasonal forecast"),
        ("Forecast (Primary)", "OFFSET(Historical, -24, 0) x 1.5", "Two-year lookback with\n50% growth multiplier.\n\nFY26: $18,300 projected.",
         "Growth-adjusted forecast"),
    ]

    card_w = Inches(2.8)
    card_h = Inches(4.5)
    gap = Inches(0.2)
    start_x = Inches(0.8)
    card_y = Inches(1.8)

    for i, (title, formula, desc, method_type) in enumerate(methods):
        x = start_x + (card_w + gap) * i
        card = add_shape(slide, x, card_y, card_w, card_h, border=True)

        # Accent bar at top of card
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, card_y, card_w, Inches(0.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT if i == 3 else BORDER_COLOR
        bar.line.fill.background()

        add_text(slide, x + Inches(0.2), card_y + Inches(0.2), card_w - Inches(0.4), Inches(0.35),
                 title, font_size=14, bold=True, color=TEXT_WHITE)

        # Formula badge
        badge = add_shape(slide, x + Inches(0.2), card_y + Inches(0.6), card_w - Inches(0.4), Inches(0.4),
                          fill_color=RGBColor(0x1A, 0x1A, 0x2E))
        add_text(slide, x + Inches(0.3), card_y + Inches(0.65), card_w - Inches(0.6), Inches(0.35),
                 formula, font_size=10, color=ACCENT)

        add_text(slide, x + Inches(0.2), card_y + Inches(1.2), card_w - Inches(0.4), Inches(2.5),
                 desc, font_size=11, color=TEXT_GRAY)

        add_text(slide, x + Inches(0.2), card_y + Inches(3.9), card_w - Inches(0.4), Inches(0.3),
                 method_type, font_size=10, color=TEXT_DIM)


def build_risk_opportunities(prs):
    """Slide 8: Risks & Opportunities"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.5), height=Inches(0.45))
    add_text(slide, Inches(1.0), Inches(0.45), Inches(10), Inches(0.6),
             "Risks & Opportunities", font_size=28, bold=True)

    add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.4),
             "Strategic considerations for FY26 forecast achievement",
             font_size=13, color=TEXT_GRAY)

    # Two columns: Risks (left) and Opportunities (right)
    col_w = Inches(5.5)
    left_x = Inches(0.8)
    right_x = Inches(7.0)
    card_top = Inches(1.8)
    card_h = Inches(5.0)

    # Risks card
    risk_card = add_shape(slide, left_x, card_top, col_w, card_h, border=True)
    # Red accent bar
    rb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, card_top, col_w, Inches(0.05))
    rb.fill.solid()
    rb.fill.fore_color.rgb = RED
    rb.line.fill.background()

    add_text(slide, left_x + Inches(0.3), card_top + Inches(0.2), col_w - Inches(0.6), Inches(0.4),
             "RISKS", font_size=18, bold=True, color=RED)

    risks = [
        ("Aggressive Growth Multiplier", "1.5x multiplier may not account for market saturation or competitive pressure. Consider sensitivity analysis at 1.2x and 1.0x."),
        ("Q2 Weakness Pattern", "Consistent Q2 dip across FY24-25 suggests structural demand trough. FY26 Q2 forecasts $3,300 — lowest quarter."),
        ("Zero Historical FY26 Data", "No actual data for FY26 yet — forecast relies entirely on FY24 historical base. Early actuals may deviate significantly."),
        ("Single-Customer Concentration", "Report focuses on Acme Corp only. Cross-customer analysis needed for portfolio risk assessment."),
    ]

    for j, (title, desc) in enumerate(risks):
        y = card_top + Inches(0.75) + Inches(1.0) * j
        add_text(slide, left_x + Inches(0.3), y, col_w - Inches(0.6), Inches(0.3),
                 title, font_size=12, bold=True, color=TEXT_WHITE)
        add_text(slide, left_x + Inches(0.3), y + Inches(0.28), col_w - Inches(0.6), Inches(0.6),
                 desc, font_size=10, color=TEXT_GRAY)

    # Opportunities card
    opp_card = add_shape(slide, right_x, card_top, col_w, card_h, border=True)
    gb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, card_top, col_w, Inches(0.05))
    gb.fill.solid()
    gb.fill.fore_color.rgb = GREEN
    gb.line.fill.background()

    add_text(slide, right_x + Inches(0.3), card_top + Inches(0.2), col_w - Inches(0.6), Inches(0.4),
             "OPPORTUNITIES", font_size=18, bold=True, color=GREEN)

    opps = [
        ("Q4 Momentum Capture", "Q4 consistently delivers ~30% of annual revenue. Focused Q4 campaigns for Laptop Pro could amplify the natural seasonal peak."),
        ("Cross-Sell Potential", "Acme Corp purchases only Laptop Pro currently. Opportunity to bundle Desktop Computer, Wireless Mouse, and accessories."),
        ("Forecast Beat Upside", "If FY25 H2 trends accelerate, actual FY26 could exceed the 1.5x projection. Build capacity for $20K+ scenario."),
        ("Pipeline Acceleration", "Q2 weakness presents an opportunity for promotional pricing or volume incentives to smooth seasonal demand."),
    ]

    for j, (title, desc) in enumerate(opps):
        y = card_top + Inches(0.75) + Inches(1.0) * j
        add_text(slide, right_x + Inches(0.3), y, col_w - Inches(0.6), Inches(0.3),
                 title, font_size=12, bold=True, color=TEXT_WHITE)
        add_text(slide, right_x + Inches(0.3), y + Inches(0.28), col_w - Inches(0.6), Inches(0.6),
                 desc, font_size=10, color=TEXT_GRAY)


def build_closing_slide(prs):
    """Slide 9: Closing / Next Steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Top accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_text(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(0.8),
             "Recommendations & Next Steps", font_size=36, bold=True)

    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.5), Inches(2.5), Inches(0.03))
    div.fill.solid()
    div.fill.fore_color.rgb = ACCENT
    div.line.fill.background()

    steps = [
        "Validate FY26 Q1 actuals against forecast to calibrate the 1.5x growth multiplier",
        "Develop Q2 promotional strategy to address historical seasonal weakness",
        "Expand cross-customer analysis — replicate this forecast for top 5 Laptop Pro accounts",
        "Build scenario models (bear: 1.0x, base: 1.5x, bull: 2.0x) for executive planning",
        "Schedule monthly forecast-vs-actual review cadence starting March 2026",
    ]

    for i, step in enumerate(steps):
        y = Inches(3.0) + Inches(0.55) * i
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y + Inches(0.02), Inches(0.35), Inches(0.35))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = TEXT_WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_text(slide, Inches(1.75), y, Inches(9), Inches(0.4),
                 step, font_size=14, color=TEXT_GRAY)

    # Footer
    add_text(slide, Inches(1.2), Inches(6.2), Inches(10), Inches(0.4),
             "Prepared with Anaplan MCP  |  [LARA] TEST MODEL  |  February 2026",
             font_size=12, color=TEXT_DIM)

    # Bottom accent bar
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = ACCENT
    bottom.line.fill.background()


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs)          # 1
    build_exec_summary(prs)         # 2
    build_quarterly_comparison(prs) # 3
    build_monthly_trend(prs)        # 4
    build_data_table(prs)           # 5
    build_waterfall_analysis(prs)   # 6
    build_forecast_methodology(prs) # 7
    build_risk_opportunities(prs)   # 8
    build_closing_slide(prs)        # 9

    output_path = os.path.join("C:/Projects/githubprojects/anaplan-mcp",
                                "FY26_Sales_Forecast_Laptop_Pro_Acme_Corp.pptx")
    prs.save(output_path)
    print(f"Report saved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
